import os
from typing import Tuple, Any
from pathlib import Path

_ML_CAPTION_PIPE = None
_TESSERACT_CHECKED = False
_TESSERACT_AVAILABLE = False


def _tesseract_available() -> bool:
    global _TESSERACT_CHECKED, _TESSERACT_AVAILABLE
    if not _TESSERACT_CHECKED:
        _TESSERACT_CHECKED = True
        try:
            import pytesseract
            pytesseract.get_tesseract_version()
            _TESSERACT_AVAILABLE = True
        except Exception:
            _TESSERACT_AVAILABLE = False
    return _TESSERACT_AVAILABLE


def _get_caption_pipe():
    global _ML_CAPTION_PIPE
    if _ML_CAPTION_PIPE is not None:
        return _ML_CAPTION_PIPE
    try:
        from transformers import pipeline
        import torch
        _ML_CAPTION_PIPE = pipeline("image-to-text", model="nlpconnect/vit-gpt2-image-captioning")
    except Exception:
        try:
            from transformers import pipeline
            _ML_CAPTION_PIPE = pipeline("image-to-text", model="Salesforce/blip-image-captioning-base")
        except Exception:
            _ML_CAPTION_PIPE = False
    return _ML_CAPTION_PIPE

CODE_EXTENSIONS = {'.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.go', '.rs', '.rb', '.php', '.c', '.cpp', '.h', '.hpp', '.cs', '.swift', '.kt', '.r', '.sh', '.bash', '.sql', '.html', '.css', '.scss', '.xml', '.json', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.md', '.tex', '.dart', '.hs', '.lua', '.pl', '.pm', '.rkt', '.zig', '.nim', '.f', '.f90', '.f95', '.ex', '.exs', '.elm'}
IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.tif', '.ico'}


class UniversalDataReader:
    @staticmethod
    def read(source: str, **kwargs) -> Tuple[Any, str]:
        import pandas as pd

        if source.startswith(('postgresql://', 'mysql://', 'sqlite:///', 'mssql+pyodbc://')):
            try:
                from sqlalchemy import create_engine
                engine = create_engine(source)
                df = pd.read_sql(kwargs.get('query', 'SELECT * FROM information_schema.tables LIMIT 10'), engine)
                return df, f"Database Source: {source.split('@')[-1]}\nQuery executed successfully."
            except Exception as e:
                return None, f"[DB_ERROR] {str(e)}"

        ext = Path(source).suffix.lower()
        try:
            if ext == '.csv':
                chunk_size = kwargs.get('chunksize', 10000)
                if chunk_size and os.path.getsize(source) > 100 * 1024 * 1024:
                    chunks = []
                    for chunk in pd.read_csv(source, chunksize=chunk_size, **kwargs):
                        chunks.append(chunk)
                    df = pd.concat(chunks, ignore_index=True)
                else:
                    df = pd.read_csv(source, **kwargs)
            elif ext in ['.xlsx', '.xls']:
                df = pd.read_excel(source, sheet_name=kwargs.get('sheet_name', 0), **kwargs)
            elif ext == '.parquet':
                df = pd.read_parquet(source, **kwargs)
            elif ext == '.feather':
                df = pd.read_feather(source, **kwargs)
            elif ext == '.json':
                df = pd.read_json(source, orient=kwargs.get('orient', 'columns'), **kwargs)
            elif ext == '.hdf5':
                df = pd.read_hdf(source, key=kwargs.get('key'), **kwargs)
            else:
                return None, f"[UNSUPPORTED] Extension '{ext}' is not supported."

            meta = f"Dataset: {Path(source).name}\nRows: {len(df):,} | Columns: {len(df.columns)}\nDtypes:\n{df.dtypes.to_dict()}\nMissing Values:\n{df.isnull().sum().to_dict()}\nFirst 3 Rows Preview:\n{df.head(3).to_markdown()}"
            return df, meta
        except Exception as e:
            return None, f"[READ_ERROR] {str(e)}]"


class DocumentReader:
    LANG_MAP = {
        '.py': 'python', '.js': 'javascript', '.ts': 'typescript',
        '.jsx': 'jsx', '.tsx': 'tsx', '.java': 'java', '.go': 'go',
        '.rs': 'rust', '.rb': 'ruby', '.php': 'php', '.c': 'c',
        '.cpp': 'cpp', '.h': 'c', '.hpp': 'cpp', '.cs': 'csharp',
        '.swift': 'swift', '.kt': 'kotlin', '.r': 'r', '.sh': 'bash',
        '.bash': 'bash', '.sql': 'sql', '.html': 'html', '.css': 'css',
        '.scss': 'scss', '.xml': 'xml', '.json': 'json', '.yaml': 'yaml',
        '.yml': 'yaml', '.toml': 'toml', '.ini': 'ini', '.md': 'markdown',
        '.tex': 'latex', '.dart': 'dart', '.hs': 'haskell', '.lua': 'lua',
        '.pl': 'perl', '.pm': 'perl', '.zig': 'zig', '.nim': 'nimrod',
        '.f': 'fortran', '.f90': 'fortran', '.f95': 'fortran',
        '.ex': 'elixir', '.exs': 'elixir', '.elm': 'elm',
    }

    @staticmethod
    def extract_text(filepath: str) -> str:
        ext = Path(filepath).suffix.lower()
        filename = Path(filepath).name

        try:
            if ext in IMAGE_EXTENSIONS:
                return DocumentReader._process_image(filepath, filename)
            if ext == '.pdf':
                return DocumentReader._process_pdf(filepath, filename)
            if ext == '.docx':
                return DocumentReader._process_docx(filepath, filename)
            if ext in ['.xlsx', '.xls']:
                return DocumentReader._process_xlsx(filepath, filename)
            if ext == '.pptx':
                return DocumentReader._process_pptx(filepath, filename)
            if ext in CODE_EXTENSIONS:
                return DocumentReader._process_code(filepath, filename, ext)
            return DocumentReader._process_text(filepath, filename)

        except Exception as e:
            return f"[Error processing {filename}: {str(e)}]"

    @staticmethod
    def _process_image(filepath: str, filename: str) -> str:
        from PIL import Image
        import numpy as np

        img = Image.open(filepath)
        parts = [
            f"[FILE: {filename}]",
            f"[TYPE: IMAGE]",
            f"[METADATA]",
            f"  Format: {img.format or 'Unknown'}",
            f"  Dimensions: {img.size[0]}x{img.size[1]}",
            f"  Mode: {img.mode}",
            f"  File size: {os.path.getsize(filepath)} bytes",
        ]

        # Convert to RGB for analysis
        if img.mode == 'RGBA':
            img_rgb = img.convert('RGB')
        elif img.mode != 'RGB':
            img_rgb = img.convert('RGB')
        else:
            img_rgb = img

        try:
            arr = np.array(img_rgb, dtype=np.float32)
            h, w, _ = arr.shape

            # === PIXEL ANALYSIS ===
            parts.append("[PIXEL ANALYSIS]")

            # Average color (overall tone)
            avg = arr.mean(axis=(0, 1))
            r_avg, g_avg, b_avg = int(avg[0]), int(avg[1]), int(avg[2])
            parts.append(f"  Average color: RGB({r_avg}, {g_avg}, {b_avg})")
            parts.append(f"  Average brightness: {arr.mean():.0f}/255")

            # Dominant colors via simple quantization (k-means lite)
            flat = arr.reshape(-1, 3)
            quantized = (flat // 64) * 64 + 32
            unique, counts = np.unique(quantized, axis=0, return_counts=True)
            top_n = min(6, len(unique))
            top_idx = np.argsort(-counts)[:top_n]
            parts.append(f"  Dominant colors (top {top_n}):")
            for idx in top_idx:
                r, g, b = int(unique[idx][0]), int(unique[idx][1]), int(unique[idx][2])
                pct = counts[idx] / counts.sum() * 100
                parts.append(f"    RGB({r}, {g}, {b}) — {pct:.1f}% of pixels")

            # Color distribution analysis
            std = arr.std(axis=(0, 1))
            color_var = tuple(int(s) for s in std)
            parts.append(f"  Color variance (std): R={color_var[0]}, G={color_var[1]}, B={color_var[2]}")

            # Check if grayscale-like
            max_channel_diff = abs(arr[:,:,0] - arr[:,:,1]).mean() + abs(arr[:,:,1] - arr[:,:,2]).mean()
            is_gray = max_channel_diff < 15
            if is_gray:
                parts.append("  Colorfulness: Grayscale or desaturated")
            elif max_channel_diff < 40:
                parts.append("  Colorfulness: Muted / low saturation")
            elif max_channel_diff < 80:
                parts.append("  Colorfulness: Moderate saturation")
            else:
                parts.append("  Colorfulness: Vivid / high saturation")

            # Warm vs cool tone
            warmth = r_avg - b_avg
            if warmth > 30:
                parts.append("  Color temperature: Warm tones")
            elif warmth < -30:
                parts.append("  Color temperature: Cool tones")
            else:
                parts.append("  Color temperature: Neutral")

            # === IMAGE CONTENT CLASSIFICATION ===
            parts.append("[CONTENT ANALYSIS]")

            # Edge detection (simple gradient) to measure complexity
            gray = arr.mean(axis=2)
            gx = np.abs(np.diff(gray, axis=1))
            gy = np.abs(np.diff(gray, axis=0))
            edge_strength = (gx.mean() + gy.mean()) / 2
            parts.append(f"  Edge strength (complexity): {edge_strength:.1f}")

            # Classify image type based on pixel stats
            if is_gray and edge_strength < 15:
                img_type = "document scan / text page"
            elif is_gray and edge_strength < 40:
                img_type = "black & white illustration or text-heavy"
            elif color_var[0] < 40 and color_var[1] < 40 and color_var[2] < 40 and edge_strength < 50:
                img_type = "screenshot / UI / flat graphic"
            elif edge_strength > 100:
                img_type = "high-detail photograph"
            elif edge_strength > 60:
                img_type = "photograph or detailed illustration"
            elif color_var[0] > 60 or color_var[1] > 60 or color_var[2] > 60:
                img_type = "colorful graphic or photo"
            else:
                img_type = "illustration or diagram"
            parts.append(f"  Estimated type: {img_type}")

            # Quadrant-based color distribution
            parts.append("  Color distribution by region:")
            regions = {
                "Top-left":   arr[:h//2, :w//2],
                "Top-right":  arr[:h//2, w//2:],
                "Bottom-left": arr[h//2:, :w//2],
                "Bottom-right": arr[h//2:, w//2:],
                "Center":     arr[h//4:3*h//4, w//4:3*w//4],
            }
            for region_name, region_arr in regions.items():
                if region_arr.size == 0:
                    continue
                ra = region_arr.mean(axis=(0, 1))
                parts.append(f"    {region_name}: RGB({int(ra[0])}, {int(ra[1])}, {int(ra[2])})")

            # Aspect ratio classification
            aspect = w / h if h > 0 else 1
            if aspect > 1.8:
                parts.append("  Aspect ratio: Wide / landscape (panoramic)")
            elif aspect > 1.2:
                parts.append("  Aspect ratio: Landscape")
            elif aspect < 0.6:
                parts.append("  Aspect ratio: Tall / portrait (vertical)")
            elif aspect < 0.9:
                parts.append("  Aspect ratio: Portrait")
            else:
                parts.append("  Aspect ratio: Square / near-square")
        except Exception as e:
            parts.append(f"  [Pixel analysis error: {e}]")

        # === ML CAPTIONING (optional) ===
        try:
            pipe = _get_caption_pipe()
            if pipe:
                caption = pipe(img)[0]['generated_text']
                parts.append(f"[ML DESCRIPTION: {caption}]")
        except Exception:
            pass

        # === OCR TEXT EXTRACTION ===
        if _tesseract_available():
            try:
                import pytesseract
                ocr_text = pytesseract.image_to_string(img).strip()
                if ocr_text:
                    parts.append("[OCR TEXT]")
                    parts.append(ocr_text)
                else:
                    parts.append("[OCR: No text detected]")
            except Exception as e:
                parts.append(f"[OCR error: {e}]")
        else:
            parts.append("[OCR: Not available - install pytesseract + Tesseract-OCR]")

        return "\n".join(parts)

    @staticmethod
    def _process_pdf(filepath: str, filename: str) -> str:
        import pdfplumber
        parts = [f"[FILE: {filename}]", f"[TYPE: PDF]"]
        with pdfplumber.open(filepath) as pdf:
            has_text = False
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text and text.strip():
                    parts.append(f"[PAGE {i+1}]")
                    parts.append(text.strip())
                    has_text = True
            if not has_text:
                parts.append("[No text content found in PDF]")
        return "\n".join(parts)

    @staticmethod
    def _process_docx(filepath: str, filename: str) -> str:
        import docx
        doc = docx.Document(filepath)
        parts = [f"[FILE: {filename}]", f"[TYPE: DOCX]"]
        has_text = False
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
                has_text = True
        if not has_text:
            parts.append("[No text content found in document]")
        return "\n".join(parts)

    @staticmethod
    def _process_xlsx(filepath: str, filename: str) -> str:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True)
        parts = [f"[FILE: {filename}]", f"[TYPE: SPREADSHEET]"]
        has_data = False
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"[SHEET: {sheet_name}] (Rows: {ws.max_row}, Cols: {ws.max_column})")
            rows_data = list(ws.iter_rows(values_only=True))
            if rows_data:
                has_data = True
                header = [str(c) if c is not None else "" for c in rows_data[0]]
                parts.append("| " + " | ".join(header) + " |")
                parts.append("| " + " | ".join(["---"] * len(header)) + " |")
                for row in rows_data[1:51]:
                    vals = [str(c)[:120] if c is not None else "" for c in row]
                    while len(vals) < len(header):
                        vals.append("")
                    parts.append("| " + " | ".join(vals) + " |")
                if len(rows_data) > 51:
                    parts.append(f"*... and {len(rows_data) - 51} more rows*")
            parts.append("")
        if not has_data:
            parts.append("[No data found in spreadsheet]")
        return "\n".join(parts)

    @staticmethod
    def _process_pptx(filepath: str, filename: str) -> str:
        import pptx
        prs = pptx.Presentation(filepath)
        parts = [f"[FILE: {filename}]", f"[TYPE: PRESENTATION]"]
        has_text = False
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"[SLIDE {i+1}]")
                parts.extend(slide_texts)
                has_text = True
        if not has_text:
            parts.append("[No text content found in presentation]")
        return "\n".join(parts)

    @staticmethod
    def _process_code(filepath: str, filename: str, ext: str) -> str:
        try:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
        except Exception:
            return f"[Error reading {filename}]"
        lang = DocumentReader.LANG_MAP.get(ext, "")
        parts = [f"[FILE: {filename}]", f"[TYPE: CODE]"]
        if lang:
            parts.append(f"[LANGUAGE: {lang}]")
        parts.append("[CONTENT START]")
        if lang:
            parts.append(f"```{lang}")
        else:
            parts.append("```")
        parts.append(content)
        parts.append("```")
        parts.append("[CONTENT END]")
        return "\n".join(parts)

    @staticmethod
    def _process_text(filepath: str, filename: str) -> str:
        file_size = os.path.getsize(filepath)
        parts = [f"[FILE: {filename}]", f"[TYPE: TEXT]"]
        try:
            if file_size > 10 * 1024 * 1024:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = "".join(f.readlines()[:1000])
            else:
                try:
                    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                except UnicodeDecodeError:
                    for encoding in ['latin-1', 'cp1252', 'utf-16', 'iso-8859-1']:
                        try:
                            with open(filepath, 'r', encoding=encoding) as f:
                                content = f.read()
                            break
                        except UnicodeDecodeError:
                            continue
                    else:
                        return f"[Error: Could not decode {filename} with any supported encoding]"
            parts.append("[CONTENT START]")
            parts.append("```")
            parts.append(content.strip())
            parts.append("```")
            parts.append("[CONTENT END]")
        except Exception as e:
            return f"[Error reading {filename}: {str(e)}]"
        return "\n".join(parts)
